import os
import magic
from typing import Union
from io import BytesIO
from pathlib import Path
from PIL import Image
import pytesseract
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import xlrd
import logging
from fastapi import HTTPException

logger = logging.getLogger(__name__)

# Configure pytesseract
try:
    pytesseract.pytesseract.tesseract_cmd = os.getenv('TESSERACT_CMD', '/usr/bin/tesseract')
except Exception as e:
    logger.warning(f"Could not set tesseract cmd: {str(e)}")

async def process_uploaded_file(file) -> str:
    """Process uploaded file and return extracted text"""
    try:
        content = await file.read()
        file_stream = BytesIO(content)
        
        # Determine file type using python-magic
        try:
            file_type = magic.from_buffer(content, mime=True)
        except Exception as e:
            logger.error(f"Error determining file type: {str(e)}")
            raise HTTPException(
                status_code=400,
                detail="Could not determine file type"
            )
        
        if file_type == 'application/pdf':
            return _extract_text_from_pdf(file_stream)
        elif file_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/msword']:
            return _extract_text_from_docx(file_stream)
        elif file_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'application/vnd.ms-excel']:
            return _extract_text_from_excel(file_stream)
        elif file_type in ['image/jpeg', 'image/png']:
            return _extract_text_from_image(file_stream)
        elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            return _extract_text_from_pptx(file_stream)
        elif file_type == 'text/plain':
            return content.decode('utf-8')
        else:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type: {file_type}"
            )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error processing file: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Could not process file: {str(e)}"
        )

def _extract_text_from_pdf(file_stream) -> str:
    try:
        reader = PdfReader(file_stream)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text
    except Exception as e:
        logger.error(f"Error extracting PDF text: {str(e)}")
        raise

def _extract_text_from_docx(file_stream) -> str:
    try:
        doc = Document(file_stream)
        return "\n".join([para.text for para in doc.paragraphs if para.text])
    except Exception as e:
        logger.error(f"Error extracting DOCX text: {str(e)}")
        raise

def _extract_text_from_excel(file_stream) -> str:
    try:
        # Try reading with pandas first
        df = pd.read_excel(file_stream)
        return df.to_string()
    except Exception as e:
        logger.warning(f"Pandas failed, trying xlrd: {str(e)}")
        try:
            # Fallback to xlrd for older formats
            workbook = xlrd.open_workbook(file_contents=file_stream.read())
            sheets = [workbook.sheet_by_index(i) for i in range(workbook.nsheets)]
            text = ""
            for sheet in sheets:
                for row in range(sheet.nrows):
                    text += "\t".join([str(cell.value) for cell in sheet.row(row)]) + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting Excel text: {str(e)}")
            raise

def _extract_text_from_image(file_stream) -> str:
    try:
        image = Image.open(file_stream)
        return pytesseract.image_to_string(image)
    except Exception as e:
        logger.error(f"Error extracting image text: {str(e)}")
        raise

def _extract_text_from_pptx(file_stream) -> str:
    try:
        prs = Presentation(file_stream)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        logger.error(f"Error extracting PPTX text: {str(e)}")
        raise
